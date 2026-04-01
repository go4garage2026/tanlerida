"""
Tangred — Arch Grant 2026 Professional Pitch Deck Generator
Generates a polished 12-slide PPTX with brand-consistent design.

Design system:
  • Background:  Rich Black (#090909)
  • Primary:     Tangred Red (#C0392B)
  • Accent:      Heritage Gold (#BFA07A)
  • Text:        Pure White (#FFFFFF) / Soft Gray (#AAAAAA)
  • Cards:       Charcoal (#1A1A1A)
"""

import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Paths ──────────────────────────────────────────────────────────
ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
OUT_DIR = os.path.join(ROOT, "pitch_video")
os.makedirs(OUT_DIR, exist_ok=True)
OUTPUT = os.path.join(OUT_DIR, "Tangred_Arch_Grant_2026_Pitch_Deck.pptx")

# ── Brand Colors ───────────────────────────────────────────────────
BLACK = RGBColor(0x09, 0x09, 0x09)
DARK_BG = RGBColor(0x0D, 0x0D, 0x0D)
CHARCOAL = RGBColor(0x1A, 0x1A, 0x1A)
CARD_BG = RGBColor(0x1E, 0x1E, 0x1E)
RED = RGBColor(0xC0, 0x39, 0x2B)
BRIGHT_RED = RGBColor(0xE7, 0x4C, 0x3C)
GOLD = RGBColor(0xBF, 0xA0, 0x7A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xAA, 0xAA, 0xAA)
MID_GRAY = RGBColor(0x66, 0x66, 0x66)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)

# ── Presentation Setup ─────────────────────────────────────────────
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height


# ── Helper Functions ───────────────────────────────────────────────

def set_slide_bg(slide, color=BLACK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color=None, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    return shape


def add_rounded_rect(slide, left, top, width, height, fill_color=CARD_BG, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    # Adjust corner radius
    shape.adjustments[0] = 0.05
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Segoe UI", anchor=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    tf.vertical_anchor = anchor
    return txBox


def add_multiline_text(slide, left, top, width, height, lines, font_size=16,
                       color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                       line_spacing=1.3, font_name="Segoe UI"):
    """Add a text box with multiple formatted lines."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None

    for i, line_data in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        if isinstance(line_data, tuple):
            text, l_color, l_bold, l_size = line_data
        else:
            text, l_color, l_bold, l_size = line_data, color, bold, font_size

        p.text = text
        p.font.size = Pt(l_size)
        p.font.color.rgb = l_color
        p.font.bold = l_bold
        p.font.name = font_name
        p.alignment = alignment
        p.space_after = Pt(font_size * 0.4)
        if line_spacing:
            p.line_spacing = line_spacing

    return txBox


def add_accent_bar(slide, left, top, width, color=RED):
    return add_shape(slide, left, top, width, Pt(4), fill_color=color)


def add_gold_separator(slide, left, top, width=Inches(2)):
    return add_shape(slide, left, top, width, Pt(2), fill_color=GOLD)


def add_slide_number(slide, num, total=12):
    add_text_box(slide, Inches(12.2), Inches(7.05), Inches(1), Inches(0.4),
                 f"{num}/{total}", font_size=11, color=MID_GRAY,
                 alignment=PP_ALIGN.RIGHT)


def add_top_bar(slide):
    add_shape(slide, Inches(0), Inches(0), SLIDE_W, Pt(5), fill_color=RED)


def add_bottom_bar(slide):
    add_shape(slide, Inches(0), Inches(7.46), SLIDE_W, Pt(3), fill_color=GOLD)


def add_section_header(slide, title, subtitle=None, y_offset=Inches(0.6)):
    add_text_box(slide, Inches(0.8), y_offset, Inches(8), Inches(0.7),
                 title, font_size=38, color=RED, bold=True,
                 font_name="Segoe UI Semibold")
    # Underline bar
    add_gold_separator(slide, Inches(0.8), y_offset + Inches(0.65), Inches(2.5))
    if subtitle:
        add_text_box(slide, Inches(0.8), y_offset + Inches(0.85), Inches(10), Inches(0.5),
                     subtitle, font_size=20, color=GOLD, font_name="Segoe UI")


def add_bullet_card(slide, left, top, width, height, icon, title, body, accent=RED):
    card = add_rounded_rect(slide, left, top, width, height, fill_color=CARD_BG, line_color=DARK_GRAY)
    # Accent bar at top of card
    add_shape(slide, left + Pt(1), top + Pt(1), width - Pt(2), Pt(3), fill_color=accent)
    # Icon/emoji
    add_text_box(slide, left + Inches(0.2), top + Inches(0.2), Inches(0.5), Inches(0.5),
                 icon, font_size=24, color=WHITE)
    # Title
    add_text_box(slide, left + Inches(0.2), top + Inches(0.6), width - Inches(0.4), Inches(0.35),
                 title, font_size=16, color=WHITE, bold=True)
    # Body
    add_text_box(slide, left + Inches(0.2), top + Inches(0.95), width - Inches(0.4), height - Inches(1.1),
                 body, font_size=13, color=LIGHT_GRAY)
    return card


# ═══════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE / COVER
# ═══════════════════════════════════════════════════════════════════
def slide_01_title():
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_slide_bg(slide)
    add_top_bar(slide)
    add_bottom_bar(slide)

    # Large brand name
    add_text_box(slide, Inches(0.8), Inches(1.5), Inches(8), Inches(1.2),
                 "TANGRED", font_size=80, color=RED, bold=True,
                 font_name="Segoe UI Black")

    add_gold_separator(slide, Inches(0.8), Inches(2.85), Inches(3))

    # Tagline
    add_text_box(slide, Inches(0.8), Inches(3.15), Inches(8), Inches(0.6),
                 "Born in India. Built for Ambition.", font_size=30,
                 color=GOLD, font_name="Segoe UI Light")

    # Value props
    props = [
        "◆  Premium Handcrafted Indian Leather",
        "◆  AI-Powered Personal Styling (Tan Lerida™)",
        "◆  Direct-to-Consumer Platform — Live & Production-Ready",
    ]
    y = Inches(4.0)
    for prop in props:
        add_text_box(slide, Inches(0.8), y, Inches(8), Inches(0.4),
                     prop, font_size=18, color=LIGHT_GRAY)
        y += Inches(0.42)

    # Right side — screenshot placeholder
    try:
        img_path = os.path.join(ROOT, "tangred-new-project-hero.png")
        if os.path.exists(img_path):
            slide.shapes.add_picture(img_path, Inches(7.8), Inches(1.2), Inches(5), Inches(3.5))
    except Exception:
        add_rounded_rect(slide, Inches(7.8), Inches(1.2), Inches(5), Inches(3.5),
                         fill_color=CHARCOAL, line_color=GOLD)
        add_text_box(slide, Inches(8.3), Inches(2.5), Inches(4), Inches(1),
                     "[Product Screenshot]", font_size=18, color=MID_GRAY,
                     alignment=PP_ALIGN.CENTER)

    # Bottom left
    add_text_box(slide, Inches(0.8), Inches(6.5), Inches(6), Inches(0.4),
                 "Arch Grant 2026 Application  •  Tangred by Go4Garage",
                 font_size=13, color=MID_GRAY)

    # Founder
    add_text_box(slide, Inches(0.8), Inches(6.9), Inches(6), Inches(0.4),
                 "Jayti Pargal  •  Founder & CEO  •  vivek@go4garage.in",
                 font_size=13, color=MID_GRAY)

    add_slide_number(slide, 1)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 2 — THE PROBLEM
# ═══════════════════════════════════════════════════════════════════
def slide_02_problem():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_top_bar(slide)
    add_bottom_bar(slide)

    add_section_header(slide, "THE PROBLEM",
                       "India's premium leather market is broken")

    # Three problem cards
    cards = [
        ("💰", "OVERPRICED IMPORTS",
         "Luxury brands like Coach & Fossil charge 2-3x in India due to import duties, middlemen, and retail markup. A $200 bag costs $500+.",
         RED),
        ("🏭", "NO CRAFT IN DOMESTIC",
         "Indian brands compete on price, not craftsmanship. Mass-produced goods flooding the market erode consumer trust in local quality.",
         RGBColor(0xB4, 0x78, 0x28)),
        ("🤷", "ZERO PERSONALIZATION",
         "Buying premium leather online is a guessing game. No sizing guidance, no style advice, no way to visualize fit before purchase.",
         RGBColor(0x78, 0x50, 0xA0)),
    ]

    x_start = Inches(0.8)
    for i, (icon, title, body, accent) in enumerate(cards):
        add_bullet_card(slide, x_start + Inches(i * 4.0), Inches(2.5),
                        Inches(3.7), Inches(2.3), icon, title, body, accent)

    # Market gap bar
    gap_y = Inches(5.2)
    add_rounded_rect(slide, Inches(0.8), gap_y, Inches(11.7), Inches(1.0),
                     fill_color=RGBColor(0x14, 0x14, 0x14))

    stats = [
        ("₹500 – ₹2,000", "Mass Market", Inches(1.5), WHITE),
        ("◀ ─── WIDE OPEN GAP ─── ▶", "No Premium Indian DTC Brand", Inches(5.0), GOLD),
        ("₹30,000+", "Imported Luxury", Inches(10.0), WHITE),
    ]
    for val, label, x, color in stats:
        add_text_box(slide, x, gap_y + Inches(0.1), Inches(3), Inches(0.45),
                     val, font_size=22, color=color, bold=True,
                     alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x, gap_y + Inches(0.55), Inches(3), Inches(0.35),
                     label, font_size=12, color=MID_GRAY,
                     alignment=PP_ALIGN.CENTER)

    # Target customer
    add_text_box(slide, Inches(0.8), Inches(6.5), Inches(11), Inches(0.4),
                 "TARGET:  Founders, executives, and consultants in Delhi, Mumbai, Bengaluru  •  "
                 "Age 25–55  •  Income ₹15L+  •  Value craftsmanship over logos",
                 font_size=13, color=LIGHT_GRAY)

    add_slide_number(slide, 2)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 3 — THE SOLUTION
# ═══════════════════════════════════════════════════════════════════
def slide_03_solution():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_top_bar(slide)
    add_bottom_bar(slide)

    add_section_header(slide, "THE SOLUTION",
                       "Tangred + Tan Lerida™ AI Styling Assistant")

    # Left column: Tangred
    left_x = Inches(0.8)
    add_text_box(slide, left_x, Inches(2.3), Inches(5.5), Inches(0.4),
                 "TANGRED — Premium Leather, Direct to You", font_size=20,
                 color=WHITE, bold=True)

    tangred_points = [
        ("▸  Handcrafted bags, briefcases, belts, jackets & accessories", LIGHT_GRAY, False, 15),
        ("▸  Full-grain & vegetable-tanned leather, artisan production", LIGHT_GRAY, False, 15),
        ("▸  ₹1,299 – ₹24,999 — luxury quality, honest pricing", LIGHT_GRAY, False, 15),
        ("▸  Direct-to-consumer: no middlemen, 60–70% gross margins", LIGHT_GRAY, False, 15),
        ("▸  Made-to-order model minimizes inventory risk", LIGHT_GRAY, False, 15),
    ]
    add_multiline_text(slide, left_x, Inches(2.8), Inches(5.5), Inches(2.2),
                       tangred_points, line_spacing=1.5)

    # Right column: Tan Lerida
    right_x = Inches(7.0)
    add_text_box(slide, right_x, Inches(2.3), Inches(5.5), Inches(0.4),
                 "TAN LERIDA™ — Your AI Master Tailor", font_size=20,
                 color=GOLD, bold=True)

    tl_points = [
        ("▸  ₹99 paid consultation — AI-powered styling session", LIGHT_GRAY, False, 15),
        ("▸  Upload photos → AI analyzes body profile & style", LIGHT_GRAY, False, 15),
        ("▸  Share preferences → budget, occasion, need", LIGHT_GRAY, False, 15),
        ("▸  Get personalized product recommendations", LIGHT_GRAY, False, 15),
        ("▸  See AI-generated visualization of you in Tangred", LIGHT_GRAY, False, 15),
    ]
    add_multiline_text(slide, right_x, Inches(2.8), Inches(5.5), Inches(2.2),
                       tl_points, line_spacing=1.5)

    # Product screenshot
    try:
        img_path = os.path.join(ROOT, "tangred-homepage.png")
        if os.path.exists(img_path):
            slide.shapes.add_picture(img_path, Inches(0.8), Inches(5.2), Inches(5.5), Inches(1.8))
    except Exception:
        pass

    # Tan Lerida screenshot
    try:
        img_path = os.path.join(ROOT, "tangred-tan-lerida-page.png")
        if os.path.exists(img_path):
            slide.shapes.add_picture(img_path, Inches(7.0), Inches(5.2), Inches(5.5), Inches(1.8))
    except Exception:
        pass

    add_slide_number(slide, 3)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 4 — HOW TAN LERIDA WORKS
# ═══════════════════════════════════════════════════════════════════
def slide_04_how_it_works():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_top_bar(slide)
    add_bottom_bar(slide)

    add_section_header(slide, "HOW TAN LERIDA™ WORKS",
                       "4-Step AI Consultation — Like a Personal Leather Tailor on Your Phone")

    steps = [
        ("01", "UPLOAD PHOTOS", "Upload casual, formal, and full-body\nphotos. Our AI reads your silhouette,\ndrape, and personal presence.",
         "📸", RED),
        ("02", "SHARE PROFILE", "Tell us your body build, skin tone,\nstyle archetype, and budget range.\nWe understand your identity.",
         "👤", RGBColor(0xB4, 0x78, 0x28)),
        ("03", "AI ANALYSIS", "Google Gemini analyzes your photos.\nAnthropic Claude generates personalized\nrecommendations matched to you.",
         "🧠", RGBColor(0x27, 0xAE, 0x60)),
        ("04", "SEE YOURSELF", "Get your curated product shortlist\nwith AI-generated visualization\nshowing you in Tangred leather.",
         "✨", GOLD),
    ]

    for i, (num, title, desc, icon, accent) in enumerate(steps):
        x = Inches(0.6) + Inches(i * 3.15)
        y = Inches(2.5)

        # Step card
        card = add_rounded_rect(slide, x, y, Inches(2.9), Inches(3.2),
                                fill_color=CARD_BG, line_color=DARK_GRAY)
        add_shape(slide, x + Pt(1), y + Pt(1), Inches(2.9) - Pt(2), Pt(3), fill_color=accent)

        # Step number
        add_text_box(slide, x + Inches(0.2), y + Inches(0.2), Inches(0.6), Inches(0.4),
                     num, font_size=28, color=accent, bold=True,
                     font_name="Segoe UI Black")

        # Icon
        add_text_box(slide, x + Inches(2.0), y + Inches(0.15), Inches(0.6), Inches(0.5),
                     icon, font_size=28, color=WHITE)

        # Title
        add_text_box(slide, x + Inches(0.2), y + Inches(0.7), Inches(2.5), Inches(0.4),
                     title, font_size=17, color=WHITE, bold=True)

        # Description
        add_text_box(slide, x + Inches(0.2), y + Inches(1.15), Inches(2.5), Inches(1.8),
                     desc, font_size=13, color=LIGHT_GRAY)

        # Arrow between steps
        if i < 3:
            add_text_box(slide, x + Inches(2.85), y + Inches(1.2), Inches(0.4), Inches(0.4),
                         "▶", font_size=16, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

    # Bottom: AI stack summary
    add_rounded_rect(slide, Inches(0.6), Inches(6.0), Inches(12.1), Inches(0.9),
                     fill_color=RGBColor(0x14, 0x14, 0x14))

    tech_labels = [
        ("Google Gemini 1.5 Pro", "Vision Analysis"),
        ("Anthropic Claude 3.5", "Recommendations"),
        ("Pinecone Vector DB", "Semantic Search"),
        ("Replicate / SDXL", "Outfit Visualization"),
    ]
    for i, (tech, role) in enumerate(tech_labels):
        tx = Inches(0.9) + Inches(i * 3.05)
        add_text_box(slide, tx, Inches(6.1), Inches(2.8), Inches(0.35),
                     tech, font_size=14, color=WHITE, bold=True,
                     alignment=PP_ALIGN.CENTER)
        add_text_box(slide, tx, Inches(6.45), Inches(2.8), Inches(0.3),
                     role, font_size=11, color=MID_GRAY,
                     alignment=PP_ALIGN.CENTER)

    add_slide_number(slide, 4)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 5 — MARKET OPPORTUNITY
# ═══════════════════════════════════════════════════════════════════
def slide_05_market():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_top_bar(slide)
    add_bottom_bar(slide)

    add_section_header(slide, "MARKET OPPORTUNITY",
                       "Massive, Growing, and Wide Open")

    # TAM / SAM / SOM
    circles_data = [
        ("TAM", "$8.2B", "India Personal Luxury\nGoods Market (2026)", Inches(3.5), Inches(3.1), Inches(3.5), RED),
        ("SAM", "$1.4B", "Premium Leather Goods\nOnline Addressable", Inches(5.3), Inches(3.4), Inches(2.5), RGBColor(0xB4, 0x78, 0x28)),
        ("SOM", "$12M", "Year 3 Target\nDTC + AI Styling", Inches(6.6), Inches(3.8), Inches(1.7), GOLD),
    ]

    for label, value, desc, cx, cy, size, color in circles_data:
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx, cy, size, size)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0x14, 0x14, 0x14)
        shape.line.color.rgb = color
        shape.line.width = Pt(2)

        center_x = cx + Inches(0.1)
        center_y = cy + size / 2 - Inches(0.5)
        add_text_box(slide, center_x, center_y - Inches(0.25), size - Inches(0.2), Inches(0.3),
                     label, font_size=13, color=color, bold=True,
                     alignment=PP_ALIGN.CENTER)
        add_text_box(slide, center_x, center_y + Inches(0.05), size - Inches(0.2), Inches(0.35),
                     value, font_size=24, color=WHITE, bold=True,
                     alignment=PP_ALIGN.CENTER)
        add_text_box(slide, center_x, center_y + Inches(0.45), size - Inches(0.2), Inches(0.6),
                     desc, font_size=10, color=LIGHT_GRAY,
                     alignment=PP_ALIGN.CENTER)

    # Right side: key insights
    right_x = Inches(9.0)
    insights = [
        ("🚀", "Fastest Growing", "India's luxury market growing\nat 12-15% CAGR"),
        ("🎯", "Underserved Gap", "₹2K–₹30K segment has no\npremium Indian DTC player"),
        ("🌏", "Diaspora Play", "US, UK, Middle East — growing\ndemand for Indian craftsmanship"),
        ("📱", "Digital Native", "60%+ of luxury discovery\nhappens online in India"),
    ]

    for i, (icon, title, desc) in enumerate(insights):
        iy = Inches(2.3) + Inches(i * 1.2)
        add_text_box(slide, right_x, iy, Inches(0.5), Inches(0.4),
                     icon, font_size=18, color=WHITE)
        add_text_box(slide, right_x + Inches(0.5), iy, Inches(3.5), Inches(0.3),
                     title, font_size=15, color=WHITE, bold=True)
        add_text_box(slide, right_x + Inches(0.5), iy + Inches(0.3), Inches(3.5), Inches(0.5),
                     desc, font_size=12, color=LIGHT_GRAY)

    add_slide_number(slide, 5)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 6 — BUSINESS MODEL
# ═══════════════════════════════════════════════════════════════════
def slide_06_business_model():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_top_bar(slide)
    add_bottom_bar(slide)

    add_section_header(slide, "BUSINESS MODEL",
                       "Two Revenue Engines — Products + AI Consultations")

    # Revenue Stream 1: Products
    add_rounded_rect(slide, Inches(0.8), Inches(2.3), Inches(5.5), Inches(3.0),
                     fill_color=CARD_BG, line_color=DARK_GRAY)
    add_shape(slide, Inches(0.8), Inches(2.3), Inches(5.5), Pt(4), fill_color=RED)

    add_text_box(slide, Inches(1.1), Inches(2.5), Inches(5), Inches(0.4),
                 "💼  PRIMARY — Product Sales", font_size=20, color=RED, bold=True)

    product_lines = [
        ("Price Range:", "₹1,299 – ₹24,999", WHITE),
        ("Average Order Value:", "₹8,999", WHITE),
        ("Gross Margin:", "60–70%", RGBColor(0x2E, 0xCC, 0x71)),
        ("Model:", "Made-to-order (low inventory risk)", LIGHT_GRAY),
        ("Categories:", "Bags, Briefcases, Belts, Jackets, Accessories", LIGHT_GRAY),
        ("Distribution:", "100% DTC — zero retail overhead", LIGHT_GRAY),
    ]
    y = Inches(3.1)
    for label, value, vcolor in product_lines:
        add_text_box(slide, Inches(1.1), y, Inches(2.2), Inches(0.3),
                     label, font_size=13, color=MID_GRAY)
        add_text_box(slide, Inches(3.3), y, Inches(2.8), Inches(0.3),
                     value, font_size=13, color=vcolor, bold=True)
        y += Inches(0.35)

    # Revenue Stream 2: Tan Lerida
    add_rounded_rect(slide, Inches(7.0), Inches(2.3), Inches(5.5), Inches(3.0),
                     fill_color=CARD_BG, line_color=DARK_GRAY)
    add_shape(slide, Inches(7.0), Inches(2.3), Inches(5.5), Pt(4), fill_color=GOLD)

    add_text_box(slide, Inches(7.3), Inches(2.5), Inches(5), Inches(0.4),
                 "🧠  SECONDARY — Tan Lerida™ Consultations", font_size=20,
                 color=GOLD, bold=True)

    tl_lines = [
        ("Price:", "₹99 + GST per session", WHITE),
        ("Net Revenue:", "~₹80/session after fees", WHITE),
        ("Conversion to Purchase:", "Estimated 10–15%", RGBColor(0x2E, 0xCC, 0x71)),
        ("Purpose:", "Acquisition + engagement funnel", LIGHT_GRAY),
        ("Data Moat:", "Body profiles + style preferences", LIGHT_GRAY),
        ("Repeat Usage:", "Unlimited sessions per user", LIGHT_GRAY),
    ]
    y = Inches(3.1)
    for label, value, vcolor in tl_lines:
        add_text_box(slide, Inches(7.3), y, Inches(2.2), Inches(0.3),
                     label, font_size=13, color=MID_GRAY)
        add_text_box(slide, Inches(9.5), y, Inches(2.8), Inches(0.3),
                     value, font_size=13, color=vcolor, bold=True)
        y += Inches(0.35)

    # Unit economics summary
    add_rounded_rect(slide, Inches(0.8), Inches(5.7), Inches(11.7), Inches(1.0),
                     fill_color=RGBColor(0x14, 0x14, 0x14))

    metrics = [
        ("₹99", "Entry Price\n(Tan Lerida)", GOLD),
        ("→", "", MID_GRAY),
        ("₹8,999", "Avg. Order\nValue", WHITE),
        ("→", "", MID_GRAY),
        ("60–70%", "Gross\nMargin", RGBColor(0x2E, 0xCC, 0x71)),
        ("→", "", MID_GRAY),
        ("91x", "Upsell\nMultiplier", RED),
    ]
    mx = Inches(1.2)
    for val, label, color in metrics:
        if val == "→":
            add_text_box(slide, mx, Inches(5.9), Inches(0.5), Inches(0.5),
                         "→", font_size=22, color=MID_GRAY,
                         alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            mx += Inches(0.5)
        else:
            add_text_box(slide, mx, Inches(5.8), Inches(1.8), Inches(0.45),
                         val, font_size=24, color=color, bold=True,
                         alignment=PP_ALIGN.CENTER)
            add_text_box(slide, mx, Inches(6.25), Inches(1.8), Inches(0.4),
                         label, font_size=10, color=MID_GRAY,
                         alignment=PP_ALIGN.CENTER)
            mx += Inches(1.8)

    add_slide_number(slide, 6)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 7 — COMPETITIVE LANDSCAPE
# ═══════════════════════════════════════════════════════════════════
def slide_07_competition():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_top_bar(slide)
    add_bottom_bar(slide)

    add_section_header(slide, "COMPETITIVE LANDSCAPE",
                       "No One Combines Craft + DTC + AI")

    # Table headers
    headers = ["BRAND", "CRAFT QUALITY", "DTC / ONLINE", "AI STYLING", "PRICE RANGE"]
    col_widths = [Inches(2.5), Inches(2.2), Inches(2.2), Inches(2.2), Inches(2.2)]
    x_start = Inches(0.8)
    y_header = Inches(2.5)

    # Header row
    add_rounded_rect(slide, x_start, y_header, Inches(11.3), Inches(0.5),
                     fill_color=RGBColor(0x20, 0x20, 0x20))
    cx = x_start
    for i, header in enumerate(headers):
        add_text_box(slide, cx + Inches(0.15), y_header + Inches(0.05),
                     col_widths[i], Inches(0.4),
                     header, font_size=12, color=GOLD, bold=True)
        cx += col_widths[i]

    # Data rows
    rows = [
        ("Hidesign", "Mass-produced", "Retail-first", "✗  None", "₹2K–₹15K", False),
        ("Nappa Dori", "Artisanal", "Limited online", "✗  None", "₹5K–₹30K", False),
        ("Da Milano", "Premium", "Mall retail", "✗  None", "₹5K–₹25K", False),
        ("Coach / Fossil", "Varies", "No India DTC", "✗  None", "₹15K–₹80K", False),
        ("TANGRED", "Handcrafted ✓", "100% DTC ✓", "✓  Tan Lerida™", "₹1.3K–₹25K", True),
    ]

    for r, (brand, craft, dtc, ai, price, highlight) in enumerate(rows):
        ry = y_header + Inches(0.6) + Inches(r * 0.55)
        bg = RGBColor(0x1A, 0x0D, 0x0B) if highlight else (CARD_BG if r % 2 == 0 else RGBColor(0x16, 0x16, 0x16))
        border = RED if highlight else None

        add_rounded_rect(slide, x_start, ry, Inches(11.3), Inches(0.5),
                         fill_color=bg, line_color=border)

        vals = [brand, craft, dtc, ai, price]
        cx = x_start
        for i, val in enumerate(vals):
            color = WHITE if highlight else LIGHT_GRAY
            if highlight and "✓" in val:
                color = RGBColor(0x2E, 0xCC, 0x71)
            elif "✗" in val:
                color = RGBColor(0x88, 0x44, 0x44)
            bld = highlight
            add_text_box(slide, cx + Inches(0.15), ry + Inches(0.05),
                         col_widths[i], Inches(0.4),
                         val, font_size=13, color=color, bold=bld)
            cx += col_widths[i]

    # Moat callout
    add_rounded_rect(slide, Inches(0.8), Inches(5.8), Inches(11.3), Inches(1.0),
                     fill_color=RGBColor(0x14, 0x14, 0x14), line_color=GOLD)
    add_text_box(slide, Inches(1.2), Inches(5.9), Inches(10.5), Inches(0.4),
                 "🏰  OUR MOAT", font_size=18, color=GOLD, bold=True)
    add_text_box(slide, Inches(1.2), Inches(6.3), Inches(10.5), Inches(0.4),
                 "Every Tan Lerida consultation captures body profiles, style preferences, and purchase intent — "
                 "proprietary demand data that tells us exactly what to build next.",
                 font_size=14, color=LIGHT_GRAY)

    add_slide_number(slide, 7)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 8 — TECHNOLOGY STACK
# ═══════════════════════════════════════════════════════════════════
def slide_08_technology():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_top_bar(slide)
    add_bottom_bar(slide)

    add_section_header(slide, "TECHNOLOGY & PRODUCT",
                       "Production-Ready Full-Stack Platform — Live Today")

    # Architecture blocks
    layers = [
        ("FRONTEND", [
            "Next.js 16 + React 19",
            "TypeScript 5 + Tailwind CSS 4",
            "Framer Motion Animations",
            "Zustand State Management",
        ], RED, Inches(0.8)),
        ("BACKEND & DATA", [
            "Next.js API Routes (25+)",
            "Prisma ORM + PostgreSQL",
            "NextAuth.js v5 (OAuth + Email)",
            "Razorpay Payments",
        ], RGBColor(0x27, 0xAE, 0x60), Inches(4.2)),
        ("AI / ML PIPELINE", [
            "Google Gemini 1.5 Pro (Vision)",
            "Anthropic Claude 3.5 (Recs)",
            "Pinecone Vector Search",
            "Replicate Image Generation",
        ], GOLD, Inches(7.6)),
        ("INFRA & OPS", [
            "Railway.app Deployment",
            "Docker Containerization",
            "Cloudinary CDN",
            "Resend Email Service",
        ], RGBColor(0x34, 0x98, 0xDB), Inches(11.0)),
    ]

    for title, items, accent, x in layers:
        y = Inches(2.4)
        card = add_rounded_rect(slide, x, y, Inches(3.1), Inches(2.8),
                                fill_color=CARD_BG, line_color=DARK_GRAY)
        add_shape(slide, x + Pt(1), y + Pt(1), Inches(3.1) - Pt(2), Pt(3), fill_color=accent)

        add_text_box(slide, x + Inches(0.2), y + Inches(0.2), Inches(2.7), Inches(0.35),
                     title, font_size=15, color=accent, bold=True)

        iy = y + Inches(0.65)
        for item in items:
            add_text_box(slide, x + Inches(0.2), iy, Inches(2.7), Inches(0.3),
                         f"▸ {item}", font_size=12, color=LIGHT_GRAY)
            iy += Inches(0.38)

    # Stats bar
    stats_y = Inches(5.6)
    add_rounded_rect(slide, Inches(0.8), stats_y, Inches(11.7), Inches(1.2),
                     fill_color=RGBColor(0x14, 0x14, 0x14))

    stats = [
        ("200+", "Files"),
        ("15,000+", "Lines of Code"),
        ("50+", "Components"),
        ("25+", "API Routes"),
        ("15+", "DB Models"),
        ("v1.0", "Production Ready"),
    ]
    for i, (val, label) in enumerate(stats):
        sx = Inches(1.0) + Inches(i * 2.0)
        add_text_box(slide, sx, stats_y + Inches(0.15), Inches(1.8), Inches(0.5),
                     val, font_size=28, color=WHITE, bold=True,
                     alignment=PP_ALIGN.CENTER)
        add_text_box(slide, sx, stats_y + Inches(0.65), Inches(1.8), Inches(0.3),
                     label, font_size=11, color=MID_GRAY,
                     alignment=PP_ALIGN.CENTER)

    add_slide_number(slide, 8)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 9 — TRACTION & MILESTONES
# ═══════════════════════════════════════════════════════════════════
def slide_09_traction():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_top_bar(slide)
    add_bottom_bar(slide)

    add_section_header(slide, "TRACTION & MILESTONES",
                       "Built, Shipped, and Ready to Scale")

    # Timeline
    milestones = [
        ("Q1 2026", "Platform Built", "Full-stack e-commerce +\nTan Lerida AI — production ready", RED, True),
        ("Q2 2026", "Arch Grant", "Relocate to St. Louis,\nlaunch beta to early adopters", GOLD, False),
        ("Q3 2026", "First Revenue", "Onboard artisan partners,\nfirst 100 paying customers", RGBColor(0x27, 0xAE, 0x60), False),
        ("Q4 2026", "Scale", "1,000+ users, expand catalog,\nmobile app development", RGBColor(0x34, 0x98, 0xDB), False),
    ]

    for i, (period, title, desc, accent, done) in enumerate(milestones):
        x = Inches(0.6) + Inches(i * 3.15)
        y = Inches(2.5)

        card = add_rounded_rect(slide, x, y, Inches(2.9), Inches(2.2),
                                fill_color=CARD_BG, line_color=DARK_GRAY)
        add_shape(slide, x + Pt(1), y + Pt(1), Inches(2.9) - Pt(2), Pt(3), fill_color=accent)

        # Status badge
        badge = "✓ DONE" if done else "PLANNED"
        badge_color = RGBColor(0x2E, 0xCC, 0x71) if done else MID_GRAY
        add_text_box(slide, x + Inches(1.6), y + Inches(0.15), Inches(1.2), Inches(0.3),
                     badge, font_size=10, color=badge_color, bold=True,
                     alignment=PP_ALIGN.RIGHT)

        add_text_box(slide, x + Inches(0.2), y + Inches(0.15), Inches(1.5), Inches(0.3),
                     period, font_size=14, color=accent, bold=True)
        add_text_box(slide, x + Inches(0.2), y + Inches(0.5), Inches(2.5), Inches(0.35),
                     title, font_size=17, color=WHITE, bold=True)
        add_text_box(slide, x + Inches(0.2), y + Inches(0.95), Inches(2.5), Inches(1.0),
                     desc, font_size=13, color=LIGHT_GRAY)

        # Connecting arrow
        if i < 3:
            add_text_box(slide, x + Inches(2.85), y + Inches(0.85), Inches(0.4), Inches(0.4),
                         "▶", font_size=16, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

    # What's already built
    built_y = Inches(5.1)
    add_text_box(slide, Inches(0.8), built_y, Inches(3), Inches(0.4),
                 "✅  ALREADY DELIVERED", font_size=16, color=RGBColor(0x2E, 0xCC, 0x71), bold=True)

    delivered = [
        "Full e-commerce: catalog, cart, checkout, orders, wishlist",
        "Tan Lerida AI: vision analysis, recommendations, visualization",
        "Auth: Email + Google OAuth with email verification",
        "Payments: Razorpay integration (products + consultations)",
        "Database: 15+ PostgreSQL models with Prisma ORM",
        "Deployed: Railway.app with Docker containerization",
    ]
    for i, item in enumerate(delivered):
        col = 0 if i < 3 else 1
        row = i % 3
        x = Inches(0.8) + Inches(col * 6.2)
        y = built_y + Inches(0.4) + Inches(row * 0.35)
        add_text_box(slide, x, y, Inches(6), Inches(0.3),
                     f"▸  {item}", font_size=12, color=LIGHT_GRAY)

    add_slide_number(slide, 9)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 10 — WHY ST. LOUIS / ARCH GRANT FIT
# ═══════════════════════════════════════════════════════════════════
def slide_10_stl():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_top_bar(slide)
    add_bottom_bar(slide)

    add_section_header(slide, "WHY ST. LOUIS",
                       "Tangred's U.S. Headquarters — Committed to the STL Ecosystem")

    cards = [
        ("🏢", "GATEWAY TO THE U.S.",
         "St. Louis is the ideal launchpad for our U.S. expansion. "
         "Lower cost of living than coastal cities means our $75K grant stretches further. "
         "Central location for shipping nationwide.",
         RED),
        ("🤝", "ECOSYSTEM FIT",
         "Arch Grant's mentor network, Cortex innovation district, and Wash U connections "
         "provide exactly the support an international founder needs to establish U.S. operations.",
         GOLD),
        ("🌍", "DIASPORA BRIDGE",
         "St. Louis has a growing South Asian community and strong ties to global markets. "
         "Tangred bridges Indian craftsmanship to American consumers — St. Louis is the crossroads.",
         RGBColor(0x27, 0xAE, 0x60)),
    ]

    for i, (icon, title, body, accent) in enumerate(cards):
        x = Inches(0.8) + Inches(i * 4.1)
        add_bullet_card(slide, x, Inches(2.5), Inches(3.8), Inches(2.8),
                        icon, title, body, accent)

    # Commitment section
    commit_y = Inches(5.7)
    add_rounded_rect(slide, Inches(0.8), commit_y, Inches(11.7), Inches(1.1),
                     fill_color=RGBColor(0x14, 0x14, 0x14), line_color=GOLD)

    add_text_box(slide, Inches(1.2), commit_y + Inches(0.1), Inches(10), Inches(0.4),
                 "OUR COMMITMENT TO ST. LOUIS", font_size=16, color=GOLD, bold=True)

    commitments = [
        "✓ Full-time founder relocation for 12+ months",
        "✓ Hire local talent for operations & marketing",
        "✓ Engage STL artisan community for U.S. production",
        "✓ Contribute to Arch Grant alumni network",
    ]
    cx = Inches(1.2)
    for i, c in enumerate(commitments):
        col_x = cx + Inches(i * 2.9)
        add_text_box(slide, col_x, commit_y + Inches(0.5), Inches(2.8), Inches(0.4),
                     c, font_size=12, color=LIGHT_GRAY)

    add_slide_number(slide, 10)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 11 — THE ASK (USE OF FUNDS)
# ═══════════════════════════════════════════════════════════════════
def slide_11_ask():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_top_bar(slide)
    add_bottom_bar(slide)

    add_section_header(slide, "THE ASK",
                       "$75,000 Equity-Free Grant — Here's How We'll Use It")

    # Fund allocation
    allocations = [
        ("40%", "$30,000", "Product & Inventory",
         "Scale artisan production, source premium leather, build initial inventory for U.S. market launch.",
         RED),
        ("25%", "$18,750", "Technology & AI",
         "Cloud infrastructure, AI API costs, mobile app development, vector database scaling.",
         GOLD),
        ("20%", "$15,000", "Marketing & Growth",
         "Digital marketing, influencer partnerships, content creation, SEO for U.S. market entry.",
         RGBColor(0x27, 0xAE, 0x60)),
        ("15%", "$11,250", "Operations & Legal",
         "U.S. entity setup, compliance, St. Louis office space, shipping logistics infrastructure.",
         RGBColor(0x34, 0x98, 0xDB)),
    ]

    for i, (pct, amount, title, desc, accent) in enumerate(allocations):
        x = Inches(0.6) + Inches(i * 3.15)
        y = Inches(2.5)

        card = add_rounded_rect(slide, x, y, Inches(2.9), Inches(3.0),
                                fill_color=CARD_BG, line_color=DARK_GRAY)
        add_shape(slide, x + Pt(1), y + Pt(1), Inches(2.9) - Pt(2), Pt(3), fill_color=accent)

        # Percentage
        add_text_box(slide, x + Inches(0.2), y + Inches(0.2), Inches(2.5), Inches(0.6),
                     pct, font_size=40, color=accent, bold=True,
                     font_name="Segoe UI Black")
        # Amount
        add_text_box(slide, x + Inches(0.2), y + Inches(0.75), Inches(2.5), Inches(0.35),
                     amount, font_size=18, color=WHITE, bold=True)
        # Title
        add_text_box(slide, x + Inches(0.2), y + Inches(1.15), Inches(2.5), Inches(0.35),
                     title, font_size=15, color=accent, bold=True)
        # Description
        add_text_box(slide, x + Inches(0.2), y + Inches(1.55), Inches(2.5), Inches(1.2),
                     desc, font_size=12, color=LIGHT_GRAY)

    # 12-month targets
    targets_y = Inches(5.9)
    add_rounded_rect(slide, Inches(0.8), targets_y, Inches(11.7), Inches(0.9),
                     fill_color=RGBColor(0x14, 0x14, 0x14))

    add_text_box(slide, Inches(1.0), targets_y + Inches(0.05), Inches(3), Inches(0.3),
                 "12-MONTH TARGETS", font_size=13, color=GOLD, bold=True)

    targets = [
        ("1,000+", "Active Users"),
        ("$150K+", "Revenue"),
        ("3–5", "Jobs Created in STL"),
        ("50+", "Artisan Partners"),
    ]
    for i, (val, label) in enumerate(targets):
        tx = Inches(1.0) + Inches(i * 3.0)
        add_text_box(slide, tx, targets_y + Inches(0.35), Inches(2.5), Inches(0.3),
                     val, font_size=22, color=WHITE, bold=True)
        add_text_box(slide, tx, targets_y + Inches(0.6), Inches(2.5), Inches(0.25),
                     label, font_size=11, color=MID_GRAY)

    add_slide_number(slide, 11)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 12 — CLOSING / THANK YOU
# ═══════════════════════════════════════════════════════════════════
def slide_12_closing():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_top_bar(slide)
    add_bottom_bar(slide)

    # Brand
    add_text_box(slide, Inches(0.8), Inches(1.4), Inches(11.7), Inches(1.0),
                 "TANGRED", font_size=72, color=RED, bold=True,
                 font_name="Segoe UI Black", alignment=PP_ALIGN.CENTER)

    add_gold_separator(slide, Inches(5.5), Inches(2.55), Inches(2.3))

    # Tagline
    add_text_box(slide, Inches(2), Inches(2.85), Inches(9.3), Inches(0.6),
                 "Crafted for Those Who Command Respect",
                 font_size=26, color=GOLD, alignment=PP_ALIGN.CENTER,
                 font_name="Segoe UI Light")

    # Key points
    points = [
        "Premium Indian Leather  •  AI-Powered Styling  •  Direct to Consumer",
        "Production-Ready Platform  •  Proven Technology  •  Ready to Scale from St. Louis",
    ]
    y = Inches(3.7)
    for pt in points:
        add_text_box(slide, Inches(1.5), y, Inches(10.3), Inches(0.4),
                     pt, font_size=17, color=LIGHT_GRAY,
                     alignment=PP_ALIGN.CENTER)
        y += Inches(0.45)

    # Contact card
    add_rounded_rect(slide, Inches(3.5), Inches(4.8), Inches(6.3), Inches(1.8),
                     fill_color=CARD_BG, line_color=GOLD)

    add_text_box(slide, Inches(3.8), Inches(4.95), Inches(5.7), Inches(0.4),
                 "LET'S BUILD TOGETHER", font_size=18, color=GOLD, bold=True,
                 alignment=PP_ALIGN.CENTER)

    contact_lines = [
        ("Jayti Pargal  •  Founder & CEO", WHITE, True, 16),
        ("vivek@go4garage.in", LIGHT_GRAY, False, 14),
    ]
    add_multiline_text(slide, Inches(3.8), Inches(5.4), Inches(5.7), Inches(1.0),
                       contact_lines, alignment=PP_ALIGN.CENTER, line_spacing=1.4)

    # Bottom
    add_text_box(slide, Inches(2), Inches(6.85), Inches(9.3), Inches(0.4),
                 "Thank you for your time and consideration.  •  Arch Grant 2026",
                 font_size=13, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

    add_slide_number(slide, 12)


# ═══════════════════════════════════════════════════════════════════
# BUILD THE DECK
# ═══════════════════════════════════════════════════════════════════
print("🏗️  Building Tangred Arch Grant 2026 Pitch Deck...")
print("=" * 55)

builders = [
    ("Slide  1/12 — Title / Cover", slide_01_title),
    ("Slide  2/12 — The Problem", slide_02_problem),
    ("Slide  3/12 — The Solution", slide_03_solution),
    ("Slide  4/12 — How Tan Lerida Works", slide_04_how_it_works),
    ("Slide  5/12 — Market Opportunity", slide_05_market),
    ("Slide  6/12 — Business Model", slide_06_business_model),
    ("Slide  7/12 — Competitive Landscape", slide_07_competition),
    ("Slide  8/12 — Technology & Product", slide_08_technology),
    ("Slide  9/12 — Traction & Milestones", slide_09_traction),
    ("Slide 10/12 — Why St. Louis", slide_10_stl),
    ("Slide 11/12 — The Ask ($75K)", slide_11_ask),
    ("Slide 12/12 — Closing / Thank You", slide_12_closing),
]

for label, builder in builders:
    print(f"  ✓ {label}")
    builder()

prs.save(OUTPUT)
print("=" * 55)
print(f"✅  Deck saved: {OUTPUT}")
print(f"📊  {len(prs.slides)} slides generated")
print("🎯  Open in PowerPoint for final review & presenter notes")
