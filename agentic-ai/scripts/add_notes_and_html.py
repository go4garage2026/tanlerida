"""
Add professional speaker notes to the Tangred Arch Grant deck,
then generate a standalone HTML preview deck for instant review.
"""

import os
from pptx import Presentation
from pptx.util import Inches

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
PPTX_PATH = os.path.join(ROOT, "pitch_video", "Tangred_Arch_Grant_2026_Pitch_Deck.pptx")
HTML_PATH = os.path.join(ROOT, "pitch_video", "Tangred_Arch_Grant_2026_Pitch_Deck.html")

# ── Speaker Notes ──────────────────────────────────────────────────
NOTES = [
    # Slide 1 — Title
    """Hi, I'm Vivek, founder of Tangred — a premium Indian leather brand with AI-powered personal styling, sold direct to consumers.

We've built a production-ready platform that combines handcrafted Indian leather goods with Tan Lerida, our AI styling assistant that gives every customer a personalized consultation.

We're applying for Arch Grant to establish our U.S. headquarters in St. Louis and scale globally.""",

    # Slide 2 — Problem
    """India's luxury leather market is broken in three ways.

First, imported brands like Coach and Fossil charge 2-3x in India due to import duties and retail markup. A $200 bag costs over $500.

Second, domestic brands compete on price, not craftsmanship. Mass-produced goods have eroded trust.

Third, buying premium leather online is impersonal. No sizing help, no style guidance, no way to see how something looks on you before buying.

India's rising professional class — founders, executives, consultants — deserves better. That's the gap we fill.""",

    # Slide 3 — Solution
    """Tangred solves this with two innovations.

First, we offer handcrafted leather goods — bags, briefcases, belts, jackets — made with full-grain and vegetable-tanned leather. Priced from $15 to $300. Direct to consumer, no middlemen, 60-70% gross margins.

Second, Tan Lerida — our AI Master Tailor. For just $1.20, customers get a full AI-powered styling consultation. Upload your photos, share your profile, and our AI recommends the exact piece that suits you. Like a personal leather tailor on your phone. No competitor offers this.""",

    # Slide 4 — How It Works
    """Here's how Tan Lerida works in four steps.

Step 1: Upload photos — casual, formal, full-body. Our AI reads your silhouette and personal presence.

Step 2: Share your profile — body build, skin tone, style preferences, budget, and occasion.

Step 3: AI Analysis — Google Gemini analyzes your photos while Anthropic Claude generates personalized recommendations.

Step 4: See Yourself — you get a curated product shortlist with an AI-generated visualization showing you wearing the recommended Tangred piece.

The entire pipeline runs on four best-in-class AI services, making it incredibly hard to replicate.""",

    # Slide 5 — Market
    """The market opportunity is massive.

India's personal luxury goods market is $8.2 billion and growing at 12-15% annually.

The addressable premium leather segment online is $1.4 billion.

Our Year 3 target is $12 million — achievable with just a tiny fraction of the market.

The key insight: there's a wide-open gap between mass-market leather under $25 and imported luxury at $400+. No premium Indian DTC brand owns this space.

Plus, the Indian diaspora in the US, UK, and Middle East represents a growing demand for authentic Indian craftsmanship.""",

    # Slide 6 — Business Model
    """We have two revenue engines.

Primary: product sales with 60-70% gross margins. Average order value is about $108. Made-to-order model means low inventory risk.

Secondary: Tan Lerida consultations at $1.20 each. Low entry price drives engagement. We estimate 10-15% conversion to purchase.

Here's the beautiful math: a $1.20 consultation leads to a $108 average order — that's a 91x upsell multiplier. And every consultation captures proprietary data about what customers actually want.""",

    # Slide 7 — Competition
    """Let me show you why no one competes with us directly.

Hidesign is mass-produced with no technology. Nappa Dori is artisanal but has limited online presence. Da Milano is stuck in malls. Coach and Fossil don't do direct-to-consumer in India.

Tangred is the only brand combining handcrafted quality, 100% DTC pricing, AND AI-powered styling.

Our moat deepens with every consultation. Each Tan Lerida session captures body profiles, style preferences, and purchase intent — proprietary demand data that tells us exactly what to build next.""",

    # Slide 8 — Technology
    """Our platform is production-ready today. Not a prototype — a deployed, working product.

The frontend runs Next.js 16 with React 19 and TypeScript. Beautiful animations with Framer Motion.

The backend has 25+ API routes, PostgreSQL with Prisma ORM, full authentication, and Razorpay payments.

The AI pipeline orchestrates four services: Gemini for vision, Claude for recommendations, Pinecone for semantic search, and Replicate for image generation.

200+ files. 15,000+ lines of code. 50+ components. This is a real product, built by a technical founder.""",

    # Slide 9 — Traction
    """Let me walk you through our timeline.

Q1 2026 — Done. We've built the complete platform. Full e-commerce, Tan Lerida AI, authentication, payments, database, deployment. All shipped.

Q2 2026 — Arch Grant. We relocate to St. Louis and launch beta to early adopters.

Q3 2026 — First revenue. We onboard artisan partners and get our first 100 paying customers.

Q4 2026 — Scale. 1,000+ users, expanded catalog, mobile app in development.

Everything on the left column is already delivered. We're not asking you to fund an idea — we're asking you to fund growth.""",

    # Slide 10 — Why STL
    """Why St. Louis? Three reasons.

First, it's the perfect gateway to the U.S. market. Lower cost of living means our grant stretches further. Central location for nationwide shipping.

Second, the ecosystem fits perfectly. Arch Grant mentors, Cortex innovation district, and Wash U connections are exactly what an international founder needs.

Third, St. Louis is a cultural bridge. The growing South Asian community and global connections make it ideal for a brand that bridges Indian craftsmanship to American consumers.

We're committed: full-time founder relocation for 12+ months, hiring local talent, and contributing to the STL startup community.""",

    # Slide 11 — The Ask
    """Here's how we'll deploy the $75,000.

40% — $30,000 for product and inventory. Scale artisan production and build initial U.S. inventory.

25% — $18,750 for technology. Cloud infrastructure, AI API costs, and mobile app development.

20% — $15,000 for marketing. Digital marketing and influencer partnerships for U.S. market entry.

15% — $11,250 for operations. U.S. entity setup, compliance, and shipping logistics.

In 12 months, we target: 1,000+ active users, $150K+ revenue, 3-5 jobs created in St. Louis, and 50+ artisan partners.""",

    # Slide 12 — Closing
    """Tangred. Born in India. Built for ambition.

Our platform is live. Our technology moat is real. Our product is handcrafted with pride.

We're ready to scale from St. Louis and bring premium Indian leather to the world — powered by AI that makes every customer feel seen.

Thank you for your time and consideration. I'd love to answer any questions.""",
]

# ── Add Notes to PPTX ─────────────────────────────────────────────
print("📝 Adding speaker notes to PPTX...")
prs = Presentation(PPTX_PATH)

for i, slide in enumerate(prs.slides):
    if i < len(NOTES):
        notes_slide = slide.notes_slide
        tf = notes_slide.notes_text_frame
        tf.text = NOTES[i]

prs.save(PPTX_PATH)
print(f"  ✓ Added notes to {len(NOTES)} slides")

# ── Generate HTML Preview Deck ─────────────────────────────────────
print("\n🌐 Generating HTML preview deck...")

HTML = """<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>TANGRED — Arch Grant 2026 Pitch Deck</title>
<style>
  @import url('https://fonts.googleapis.com/css2?family=Inter:wght@300;400;600;700;900&display=swap');

  :root {
    --bg: #090909;
    --card: #1a1a1a;
    --card2: #141414;
    --red: #C0392B;
    --bright-red: #E74C3C;
    --gold: #BFA07A;
    --white: #FFFFFF;
    --gray: #AAAAAA;
    --mid-gray: #666666;
    --dark-gray: #333333;
    --green: #2ECC71;
    --blue: #3498DB;
  }

  * { margin: 0; padding: 0; box-sizing: border-box; }

  body {
    font-family: 'Inter', 'Segoe UI', system-ui, sans-serif;
    background: #000;
    color: var(--white);
    overflow-x: hidden;
  }

  .slide {
    width: 100vw;
    min-height: 100vh;
    background: var(--bg);
    position: relative;
    display: flex;
    flex-direction: column;
    justify-content: center;
    padding: 60px 80px;
    border-top: 5px solid var(--red);
    border-bottom: 3px solid var(--gold);
  }

  .slide-number {
    position: absolute;
    bottom: 20px;
    right: 40px;
    color: var(--mid-gray);
    font-size: 14px;
  }

  h1 { font-size: 4.5rem; font-weight: 900; color: var(--red); letter-spacing: 4px; }
  h2 { font-size: 2.5rem; font-weight: 700; color: var(--red); margin-bottom: 8px; }
  .subtitle { font-size: 1.3rem; color: var(--gold); margin-bottom: 30px; font-weight: 300; }
  .gold-line { width: 160px; height: 2px; background: var(--gold); margin: 16px 0; }
  .tagline { font-size: 1.6rem; color: var(--gold); font-weight: 300; }

  .props { list-style: none; margin-top: 20px; }
  .props li { font-size: 1.1rem; color: var(--gray); padding: 6px 0; }
  .props li::before { content: '◆ '; color: var(--red); }

  .cards { display: grid; gap: 24px; margin: 30px 0; }
  .cards-3 { grid-template-columns: repeat(3, 1fr); }
  .cards-4 { grid-template-columns: repeat(4, 1fr); }
  .cards-2 { grid-template-columns: repeat(2, 1fr); }

  .card {
    background: var(--card);
    border-radius: 12px;
    padding: 28px;
    border-top: 3px solid var(--red);
    border: 1px solid var(--dark-gray);
    transition: transform 0.2s;
  }
  .card:hover { transform: translateY(-4px); }
  .card.gold-border { border-top-color: var(--gold); }
  .card.green-border { border-top-color: var(--green); }
  .card.blue-border { border-top-color: var(--blue); }
  .card.orange-border { border-top-color: #B47828; }

  .card .icon { font-size: 1.8rem; margin-bottom: 12px; }
  .card h3 { font-size: 1rem; font-weight: 700; margin-bottom: 12px; letter-spacing: 1px; }
  .card p { font-size: 0.9rem; color: var(--gray); line-height: 1.6; }
  .card .num { font-size: 2.5rem; font-weight: 900; margin-bottom: 4px; }
  .card .pct { font-size: 3rem; font-weight: 900; margin-bottom: 8px; }

  .stat-bar {
    background: var(--card2);
    border-radius: 12px;
    padding: 20px 30px;
    display: flex;
    justify-content: space-around;
    align-items: center;
    margin: 20px 0;
  }
  .stat { text-align: center; }
  .stat .val { font-size: 1.8rem; font-weight: 700; }
  .stat .label { font-size: 0.75rem; color: var(--mid-gray); margin-top: 4px; }

  table {
    width: 100%;
    border-collapse: separate;
    border-spacing: 0 6px;
    margin: 20px 0;
  }
  th {
    background: #202020;
    color: var(--gold);
    font-size: 0.8rem;
    font-weight: 700;
    letter-spacing: 1px;
    padding: 12px 16px;
    text-align: left;
  }
  td {
    background: var(--card);
    color: var(--gray);
    font-size: 0.9rem;
    padding: 12px 16px;
  }
  tr.highlight td {
    background: #1a0d0b;
    color: var(--white);
    font-weight: 600;
    border: 1px solid var(--red);
  }
  .check { color: var(--green); }
  .cross { color: #884444; }

  .moat-box {
    background: var(--card2);
    border: 1px solid var(--gold);
    border-radius: 12px;
    padding: 20px 28px;
    margin-top: 20px;
  }
  .moat-box h4 { color: var(--gold); font-size: 1.1rem; margin-bottom: 8px; }
  .moat-box p { color: var(--gray); font-size: 0.95rem; line-height: 1.5; }

  .funnel {
    display: flex;
    align-items: center;
    justify-content: center;
    gap: 12px;
    margin: 20px 0;
    flex-wrap: wrap;
  }
  .funnel .step { text-align: center; }
  .funnel .step .val { font-size: 1.6rem; font-weight: 700; }
  .funnel .step .label { font-size: 0.7rem; color: var(--mid-gray); }
  .funnel .arrow { font-size: 1.4rem; color: var(--mid-gray); }

  .footer-note {
    position: absolute;
    bottom: 30px;
    left: 80px;
    color: var(--mid-gray);
    font-size: 0.85rem;
  }

  .contact-card {
    background: var(--card);
    border: 1px solid var(--gold);
    border-radius: 16px;
    padding: 32px 48px;
    text-align: center;
    max-width: 500px;
    margin: 30px auto;
  }
  .contact-card h3 { color: var(--gold); font-size: 1.1rem; letter-spacing: 2px; margin-bottom: 16px; }
  .contact-card .name { font-size: 1.1rem; font-weight: 600; }
  .contact-card .detail { color: var(--gray); font-size: 0.95rem; margin-top: 6px; }

  .delivered { display: grid; grid-template-columns: 1fr 1fr; gap: 8px; margin-top: 16px; }
  .delivered li { color: var(--gray); font-size: 0.9rem; padding: 2px 0; list-style: none; }
  .delivered li::before { content: '▸ '; color: var(--green); }

  .commit-bar {
    background: var(--card2);
    border: 1px solid var(--gold);
    border-radius: 12px;
    padding: 16px 28px;
    display: grid;
    grid-template-columns: repeat(4, 1fr);
    gap: 16px;
    margin-top: 24px;
  }
  .commit-bar span { color: var(--gray); font-size: 0.85rem; }

  .two-col { display: grid; grid-template-columns: 1fr 1fr; gap: 40px; }
  .col-list li { color: var(--gray); font-size: 1rem; padding: 5px 0; list-style: none; }
  .col-list li::before { content: '▸ '; color: var(--gold); }

  @media (max-width: 900px) {
    .slide { padding: 30px 24px; }
    h1 { font-size: 2.8rem; }
    h2 { font-size: 1.6rem; }
    .cards-3, .cards-4 { grid-template-columns: 1fr; }
    .two-col { grid-template-columns: 1fr; }
    .commit-bar { grid-template-columns: 1fr 1fr; }
    .delivered { grid-template-columns: 1fr; }
  }
</style>
</head>
<body>

<!-- SLIDE 1: TITLE -->
<section class="slide">
  <h1>TANGRED</h1>
  <div class="gold-line"></div>
  <p class="tagline">Born in India. Built for Ambition.</p>
  <ul class="props" style="margin-top:28px">
    <li>Premium Handcrafted Indian Leather</li>
    <li>AI-Powered Personal Styling (Tan Lerida™)</li>
    <li>Direct-to-Consumer Platform — Live & Production-Ready</li>
  </ul>
  <p class="footer-note">Arch Grant 2026 Application &bull; Tangred by Go4Garage<br>Vivek &bull; Founder & CEO &bull; vivek@go4garage.in</p>
  <span class="slide-number">1/12</span>
</section>

<!-- SLIDE 2: PROBLEM -->
<section class="slide">
  <h2>THE PROBLEM</h2>
  <p class="subtitle">India's premium leather market is broken</p>
  <div class="cards cards-3">
    <div class="card">
      <div class="icon">💰</div>
      <h3 style="color:var(--red)">OVERPRICED IMPORTS</h3>
      <p>Luxury brands like Coach & Fossil charge 2-3x in India due to import duties, middlemen, and retail markup. A $200 bag costs $500+.</p>
    </div>
    <div class="card orange-border">
      <div class="icon">🏭</div>
      <h3 style="color:#B47828">NO CRAFT IN DOMESTIC</h3>
      <p>Indian brands compete on price, not craftsmanship. Mass-produced goods flooding the market erode consumer trust in local quality.</p>
    </div>
    <div class="card" style="border-top-color:#7850A0">
      <div class="icon">🤷</div>
      <h3 style="color:#7850A0">ZERO PERSONALIZATION</h3>
      <p>Buying premium leather online is a guessing game. No sizing guidance, no style advice, no way to visualize fit before purchase.</p>
    </div>
  </div>
  <div class="stat-bar">
    <div class="stat"><div class="val">₹500 – ₹2,000</div><div class="label">Mass Market</div></div>
    <div class="stat"><div class="val" style="color:var(--gold)">◀── WIDE OPEN GAP ──▶</div><div class="label">No Premium Indian DTC Brand</div></div>
    <div class="stat"><div class="val">₹30,000+</div><div class="label">Imported Luxury</div></div>
  </div>
  <p style="color:var(--gray);font-size:0.9rem;margin-top:12px">
    <strong style="color:var(--red)">TARGET:</strong> Founders, executives, consultants in Delhi, Mumbai, Bengaluru &bull; Age 25–55 &bull; Income ₹15L+ &bull; Value craftsmanship over logos
  </p>
  <span class="slide-number">2/12</span>
</section>

<!-- SLIDE 3: SOLUTION -->
<section class="slide">
  <h2>THE SOLUTION</h2>
  <p class="subtitle">Tangred + Tan Lerida™ AI Styling Assistant</p>
  <div class="two-col">
    <div>
      <h3 style="font-size:1.2rem;margin-bottom:16px;">TANGRED — Premium Leather, Direct to You</h3>
      <ul class="col-list">
        <li>Handcrafted bags, briefcases, belts, jackets & accessories</li>
        <li>Full-grain & vegetable-tanned leather, artisan production</li>
        <li>₹1,299 – ₹24,999 — luxury quality, honest pricing</li>
        <li>Direct-to-consumer: no middlemen, 60–70% gross margins</li>
        <li>Made-to-order model minimizes inventory risk</li>
      </ul>
    </div>
    <div>
      <h3 style="font-size:1.2rem;color:var(--gold);margin-bottom:16px;">TAN LERIDA™ — Your AI Master Tailor</h3>
      <ul class="col-list">
        <li>₹99 paid consultation — AI-powered styling session</li>
        <li>Upload photos → AI analyzes body profile & style</li>
        <li>Share preferences → budget, occasion, need</li>
        <li>Get personalized product recommendations</li>
        <li>See AI-generated visualization of you in Tangred</li>
      </ul>
    </div>
  </div>
  <span class="slide-number">3/12</span>
</section>

<!-- SLIDE 4: HOW IT WORKS -->
<section class="slide">
  <h2>HOW TAN LERIDA™ WORKS</h2>
  <p class="subtitle">4-Step AI Consultation — Like a Personal Leather Tailor on Your Phone</p>
  <div class="cards cards-4">
    <div class="card">
      <div class="num" style="color:var(--red)">01</div>
      <h3>📸 UPLOAD PHOTOS</h3>
      <p>Upload casual, formal, and full-body photos. Our AI reads your silhouette, drape, and personal presence.</p>
    </div>
    <div class="card orange-border">
      <div class="num" style="color:#B47828">02</div>
      <h3>👤 SHARE PROFILE</h3>
      <p>Tell us your body build, skin tone, style archetype, and budget range. We understand your identity.</p>
    </div>
    <div class="card green-border">
      <div class="num" style="color:var(--green)">03</div>
      <h3>🧠 AI ANALYSIS</h3>
      <p>Google Gemini analyzes your photos. Anthropic Claude generates personalized recommendations matched to you.</p>
    </div>
    <div class="card gold-border">
      <div class="num" style="color:var(--gold)">04</div>
      <h3>✨ SEE YOURSELF</h3>
      <p>Get your curated product shortlist with AI-generated visualization showing you in Tangred leather.</p>
    </div>
  </div>
  <div class="stat-bar">
    <div class="stat"><div class="val" style="font-size:1rem">Google Gemini 1.5 Pro</div><div class="label">Vision Analysis</div></div>
    <div class="stat"><div class="val" style="font-size:1rem">Anthropic Claude 3.5</div><div class="label">Recommendations</div></div>
    <div class="stat"><div class="val" style="font-size:1rem">Pinecone Vector DB</div><div class="label">Semantic Search</div></div>
    <div class="stat"><div class="val" style="font-size:1rem">Replicate / SDXL</div><div class="label">Outfit Visualization</div></div>
  </div>
  <span class="slide-number">4/12</span>
</section>

<!-- SLIDE 5: MARKET -->
<section class="slide">
  <h2>MARKET OPPORTUNITY</h2>
  <p class="subtitle">Massive, Growing, and Wide Open</p>
  <div class="cards cards-3" style="align-items:center">
    <div class="card" style="text-align:center;border-color:var(--red)">
      <div class="label" style="color:var(--red);font-weight:700;font-size:0.9rem;letter-spacing:2px">TAM</div>
      <div class="num" style="color:var(--white)">$8.2B</div>
      <p>India Personal Luxury Goods Market (2026)</p>
    </div>
    <div class="card orange-border" style="text-align:center">
      <div class="label" style="color:#B47828;font-weight:700;font-size:0.9rem;letter-spacing:2px">SAM</div>
      <div class="num" style="color:var(--white)">$1.4B</div>
      <p>Premium Leather Goods Online Addressable</p>
    </div>
    <div class="card gold-border" style="text-align:center">
      <div class="label" style="color:var(--gold);font-weight:700;font-size:0.9rem;letter-spacing:2px">SOM</div>
      <div class="num" style="color:var(--white)">$12M</div>
      <p>Year 3 Target — DTC + AI Styling</p>
    </div>
  </div>
  <div class="cards cards-4" style="margin-top:16px">
    <div style="display:flex;gap:12px;align-items:flex-start">
      <span style="font-size:1.4rem">🚀</span>
      <div><strong>Fastest Growing</strong><br><span style="color:var(--gray);font-size:0.9rem">India's luxury market growing at 12–15% CAGR</span></div>
    </div>
    <div style="display:flex;gap:12px;align-items:flex-start">
      <span style="font-size:1.4rem">🎯</span>
      <div><strong>Underserved Gap</strong><br><span style="color:var(--gray);font-size:0.9rem">₹2K–₹30K segment has no premium Indian DTC player</span></div>
    </div>
    <div style="display:flex;gap:12px;align-items:flex-start">
      <span style="font-size:1.4rem">🌏</span>
      <div><strong>Diaspora Play</strong><br><span style="color:var(--gray);font-size:0.9rem">US, UK, Middle East — demand for Indian craftsmanship</span></div>
    </div>
    <div style="display:flex;gap:12px;align-items:flex-start">
      <span style="font-size:1.4rem">📱</span>
      <div><strong>Digital Native</strong><br><span style="color:var(--gray);font-size:0.9rem">60%+ of luxury discovery happens online in India</span></div>
    </div>
  </div>
  <span class="slide-number">5/12</span>
</section>

<!-- SLIDE 6: BUSINESS MODEL -->
<section class="slide">
  <h2>BUSINESS MODEL</h2>
  <p class="subtitle">Two Revenue Engines — Products + AI Consultations</p>
  <div class="cards cards-2">
    <div class="card">
      <h3 style="color:var(--red);font-size:1.1rem">💼 PRIMARY — Product Sales</h3>
      <table style="border-spacing:0;margin:12px 0">
        <tr><td style="background:transparent;color:var(--mid-gray);padding:4px 0">Price Range</td><td style="background:transparent;padding:4px 0"><strong>₹1,299 – ₹24,999</strong></td></tr>
        <tr><td style="background:transparent;color:var(--mid-gray);padding:4px 0">Average Order</td><td style="background:transparent;padding:4px 0"><strong>₹8,999</strong></td></tr>
        <tr><td style="background:transparent;color:var(--mid-gray);padding:4px 0">Gross Margin</td><td style="background:transparent;color:var(--green);padding:4px 0"><strong>60–70%</strong></td></tr>
        <tr><td style="background:transparent;color:var(--mid-gray);padding:4px 0">Model</td><td style="background:transparent;color:var(--gray);padding:4px 0">Made-to-order (low risk)</td></tr>
        <tr><td style="background:transparent;color:var(--mid-gray);padding:4px 0">Distribution</td><td style="background:transparent;color:var(--gray);padding:4px 0">100% DTC — zero retail</td></tr>
      </table>
    </div>
    <div class="card gold-border">
      <h3 style="color:var(--gold);font-size:1.1rem">🧠 SECONDARY — Tan Lerida™ Consultations</h3>
      <table style="border-spacing:0;margin:12px 0">
        <tr><td style="background:transparent;color:var(--mid-gray);padding:4px 0">Price</td><td style="background:transparent;padding:4px 0"><strong>₹99 + GST</strong></td></tr>
        <tr><td style="background:transparent;color:var(--mid-gray);padding:4px 0">Net Revenue</td><td style="background:transparent;padding:4px 0"><strong>~₹80/session</strong></td></tr>
        <tr><td style="background:transparent;color:var(--mid-gray);padding:4px 0">Conversion</td><td style="background:transparent;color:var(--green);padding:4px 0"><strong>Est. 10–15%</strong></td></tr>
        <tr><td style="background:transparent;color:var(--mid-gray);padding:4px 0">Purpose</td><td style="background:transparent;color:var(--gray);padding:4px 0">Acquisition funnel</td></tr>
        <tr><td style="background:transparent;color:var(--mid-gray);padding:4px 0">Data Moat</td><td style="background:transparent;color:var(--gray);padding:4px 0">Body profiles + preferences</td></tr>
      </table>
    </div>
  </div>
  <div class="funnel" style="margin-top:24px">
    <div class="step"><div class="val" style="color:var(--gold)">₹99</div><div class="label">Entry Price</div></div>
    <div class="arrow">→</div>
    <div class="step"><div class="val">₹8,999</div><div class="label">Avg Order</div></div>
    <div class="arrow">→</div>
    <div class="step"><div class="val" style="color:var(--green)">60–70%</div><div class="label">Gross Margin</div></div>
    <div class="arrow">→</div>
    <div class="step"><div class="val" style="color:var(--red)">91x</div><div class="label">Upsell Multiplier</div></div>
  </div>
  <span class="slide-number">6/12</span>
</section>

<!-- SLIDE 7: COMPETITION -->
<section class="slide">
  <h2>COMPETITIVE LANDSCAPE</h2>
  <p class="subtitle">No One Combines Craft + DTC + AI</p>
  <table>
    <thead><tr><th>BRAND</th><th>CRAFT QUALITY</th><th>DTC / ONLINE</th><th>AI STYLING</th><th>PRICE RANGE</th></tr></thead>
    <tbody>
      <tr><td>Hidesign</td><td>Mass-produced</td><td>Retail-first</td><td class="cross">✗ None</td><td>₹2K–₹15K</td></tr>
      <tr><td>Nappa Dori</td><td>Artisanal</td><td>Limited online</td><td class="cross">✗ None</td><td>₹5K–₹30K</td></tr>
      <tr><td>Da Milano</td><td>Premium</td><td>Mall retail</td><td class="cross">✗ None</td><td>₹5K–₹25K</td></tr>
      <tr><td>Coach / Fossil</td><td>Varies</td><td>No India DTC</td><td class="cross">✗ None</td><td>₹15K–₹80K</td></tr>
      <tr class="highlight"><td><strong>TANGRED</strong></td><td class="check">Handcrafted ✓</td><td class="check">100% DTC ✓</td><td class="check">✓ Tan Lerida™</td><td>₹1.3K–₹25K</td></tr>
    </tbody>
  </table>
  <div class="moat-box">
    <h4>🏰 OUR MOAT</h4>
    <p>Every Tan Lerida consultation captures body profiles, style preferences, and purchase intent — proprietary demand data that tells us exactly what to build next.</p>
  </div>
  <span class="slide-number">7/12</span>
</section>

<!-- SLIDE 8: TECHNOLOGY -->
<section class="slide">
  <h2>TECHNOLOGY & PRODUCT</h2>
  <p class="subtitle">Production-Ready Full-Stack Platform — Live Today</p>
  <div class="cards cards-4">
    <div class="card"><h3 style="color:var(--red)">FRONTEND</h3><p>Next.js 16 + React 19<br>TypeScript 5 + Tailwind 4<br>Framer Motion<br>Zustand State Mgmt</p></div>
    <div class="card green-border"><h3 style="color:var(--green)">BACKEND & DATA</h3><p>25+ API Routes<br>Prisma ORM + PostgreSQL<br>NextAuth.js v5 (OAuth)<br>Razorpay Payments</p></div>
    <div class="card gold-border"><h3 style="color:var(--gold)">AI / ML PIPELINE</h3><p>Google Gemini 1.5 (Vision)<br>Claude 3.5 (Recommendations)<br>Pinecone Vector Search<br>Replicate Image Gen</p></div>
    <div class="card blue-border"><h3 style="color:var(--blue)">INFRA & OPS</h3><p>Railway.app Deploy<br>Docker Containerization<br>Cloudinary CDN<br>Resend Email</p></div>
  </div>
  <div class="stat-bar">
    <div class="stat"><div class="val">200+</div><div class="label">Files</div></div>
    <div class="stat"><div class="val">15,000+</div><div class="label">Lines of Code</div></div>
    <div class="stat"><div class="val">50+</div><div class="label">Components</div></div>
    <div class="stat"><div class="val">25+</div><div class="label">API Routes</div></div>
    <div class="stat"><div class="val">15+</div><div class="label">DB Models</div></div>
    <div class="stat"><div class="val" style="color:var(--green)">v1.0</div><div class="label">Production Ready</div></div>
  </div>
  <span class="slide-number">8/12</span>
</section>

<!-- SLIDE 9: TRACTION -->
<section class="slide">
  <h2>TRACTION & MILESTONES</h2>
  <p class="subtitle">Built, Shipped, and Ready to Scale</p>
  <div class="cards cards-4">
    <div class="card">
      <div style="display:flex;justify-content:space-between;align-items:center;margin-bottom:8px">
        <span style="color:var(--red);font-weight:700">Q1 2026</span>
        <span style="color:var(--green);font-size:0.75rem;font-weight:700">✓ DONE</span>
      </div>
      <h3>Platform Built</h3>
      <p>Full-stack e-commerce + Tan Lerida AI — production ready</p>
    </div>
    <div class="card gold-border">
      <div style="display:flex;justify-content:space-between;align-items:center;margin-bottom:8px">
        <span style="color:var(--gold);font-weight:700">Q2 2026</span>
        <span style="color:var(--mid-gray);font-size:0.75rem">PLANNED</span>
      </div>
      <h3>Arch Grant</h3>
      <p>Relocate to St. Louis, launch beta to early adopters</p>
    </div>
    <div class="card green-border">
      <div style="display:flex;justify-content:space-between;align-items:center;margin-bottom:8px">
        <span style="color:var(--green);font-weight:700">Q3 2026</span>
        <span style="color:var(--mid-gray);font-size:0.75rem">PLANNED</span>
      </div>
      <h3>First Revenue</h3>
      <p>Onboard artisan partners, first 100 paying customers</p>
    </div>
    <div class="card blue-border">
      <div style="display:flex;justify-content:space-between;align-items:center;margin-bottom:8px">
        <span style="color:var(--blue);font-weight:700">Q4 2026</span>
        <span style="color:var(--mid-gray);font-size:0.75rem">PLANNED</span>
      </div>
      <h3>Scale</h3>
      <p>1,000+ users, expand catalog, mobile app development</p>
    </div>
  </div>
  <h4 style="color:var(--green);margin-top:20px">✅ ALREADY DELIVERED</h4>
  <ul class="delivered">
    <li>Full e-commerce: catalog, cart, checkout, orders, wishlist</li>
    <li>Payments: Razorpay integration (products + consultations)</li>
    <li>Tan Lerida AI: vision analysis, recommendations, visualization</li>
    <li>Database: 15+ PostgreSQL models with Prisma ORM</li>
    <li>Auth: Email + Google OAuth with email verification</li>
    <li>Deployed: Railway.app with Docker containerization</li>
  </ul>
  <span class="slide-number">9/12</span>
</section>

<!-- SLIDE 10: WHY STL -->
<section class="slide">
  <h2>WHY ST. LOUIS</h2>
  <p class="subtitle">Tangred's U.S. Headquarters — Committed to the STL Ecosystem</p>
  <div class="cards cards-3">
    <div class="card">
      <div class="icon">🏢</div>
      <h3 style="color:var(--red)">GATEWAY TO THE U.S.</h3>
      <p>St. Louis is the ideal launchpad. Lower cost of living means our $75K grant stretches further. Central location for shipping nationwide.</p>
    </div>
    <div class="card gold-border">
      <div class="icon">🤝</div>
      <h3 style="color:var(--gold)">ECOSYSTEM FIT</h3>
      <p>Arch Grant mentors, Cortex innovation district, and Wash U connections provide exactly the support an international founder needs.</p>
    </div>
    <div class="card green-border">
      <div class="icon">🌍</div>
      <h3 style="color:var(--green)">DIASPORA BRIDGE</h3>
      <p>Growing South Asian community and global ties make STL ideal for bridging Indian craftsmanship to American consumers.</p>
    </div>
  </div>
  <div class="commit-bar">
    <span>✓ Full-time founder relocation 12+ months</span>
    <span>✓ Hire local talent for ops & marketing</span>
    <span>✓ Engage STL artisan community</span>
    <span>✓ Contribute to Arch Grant alumni network</span>
  </div>
  <span class="slide-number">10/12</span>
</section>

<!-- SLIDE 11: THE ASK -->
<section class="slide">
  <h2>THE ASK</h2>
  <p class="subtitle">$75,000 Equity-Free Grant — Here's How We'll Use It</p>
  <div class="cards cards-4">
    <div class="card">
      <div class="pct" style="color:var(--red)">40%</div>
      <strong>$30,000</strong>
      <h3 style="color:var(--red);margin-top:8px">Product & Inventory</h3>
      <p>Scale artisan production, source premium leather, build initial inventory for U.S. market launch.</p>
    </div>
    <div class="card gold-border">
      <div class="pct" style="color:var(--gold)">25%</div>
      <strong>$18,750</strong>
      <h3 style="color:var(--gold);margin-top:8px">Technology & AI</h3>
      <p>Cloud infrastructure, AI API costs, mobile app development, vector database scaling.</p>
    </div>
    <div class="card green-border">
      <div class="pct" style="color:var(--green)">20%</div>
      <strong>$15,000</strong>
      <h3 style="color:var(--green);margin-top:8px">Marketing & Growth</h3>
      <p>Digital marketing, influencer partnerships, content creation, SEO for U.S. market entry.</p>
    </div>
    <div class="card blue-border">
      <div class="pct" style="color:var(--blue)">15%</div>
      <strong>$11,250</strong>
      <h3 style="color:var(--blue);margin-top:8px">Operations & Legal</h3>
      <p>U.S. entity setup, compliance, STL office space, shipping logistics infrastructure.</p>
    </div>
  </div>
  <div class="stat-bar" style="margin-top:24px">
    <div class="stat"><div class="val">1,000+</div><div class="label">Active Users</div></div>
    <div class="stat"><div class="val" style="color:var(--green)">$150K+</div><div class="label">Revenue</div></div>
    <div class="stat"><div class="val">3–5</div><div class="label">Jobs in STL</div></div>
    <div class="stat"><div class="val" style="color:var(--gold)">50+</div><div class="label">Artisan Partners</div></div>
  </div>
  <p style="text-align:center;color:var(--mid-gray);font-size:0.85rem;margin-top:8px">12-MONTH TARGETS</p>
  <span class="slide-number">11/12</span>
</section>

<!-- SLIDE 12: CLOSING -->
<section class="slide" style="text-align:center;justify-content:center">
  <h1 style="font-size:5rem">TANGRED</h1>
  <div class="gold-line" style="margin:16px auto"></div>
  <p class="tagline" style="margin-bottom:24px">Crafted for Those Who Command Respect</p>
  <p style="color:var(--gray);font-size:1.05rem;margin-bottom:8px">Premium Indian Leather &bull; AI-Powered Styling &bull; Direct to Consumer</p>
  <p style="color:var(--gray);font-size:1.05rem;margin-bottom:32px">Production-Ready Platform &bull; Proven Technology &bull; Ready to Scale from St. Louis</p>
  <div class="contact-card">
    <h3>LET'S BUILD TOGETHER</h3>
    <p class="name">Vivek &bull; Founder & CEO</p>
    <p class="detail">vivek@go4garage.in</p>
    <p class="detail">github.com/G4G-EKA-Ai/TANLERIDA</p>
  </div>
  <p style="color:var(--mid-gray);font-size:0.9rem;margin-top:32px">Thank you for your time and consideration. &bull; Arch Grant 2026</p>
  <span class="slide-number">12/12</span>
</section>

</body>
</html>"""

with open(HTML_PATH, "w", encoding="utf-8") as f:
    f.write(HTML)

print(f"  ✓ Saved: {HTML_PATH}")
print("\n🎉 Pitch deck generation complete!")
print(f"   📊 PPTX: {PPTX_PATH}")
print(f"   🌐 HTML: {HTML_PATH}")
